#!/usr/bin/env python3
"""Build the PETSE review-response (rebuttal) slide deck (.pptx).

Maps each reviewer critique -> our response -> evidence -> paper location.
For showing the advisor how the TII rejection was addressed.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = '/home/jim/ros2_motion_planning_tutorials'
FIG = os.path.join(ROOT, 'figures')
OUT = os.path.join(ROOT, 'PETSE_rebuttal.pptx')

NAVY = RGBColor(0x10, 0x2A, 0x54)
BLUE = RGBColor(0x1A, 0x73, 0xE8)
RED = RGBColor(0xC6, 0x28, 0x28)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GREY = RGBColor(0x54, 0x5A, 0x66)
LIGHT = RGBColor(0xEE, 0xF2, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARKTXT = RGBColor(0x22, 0x26, 0x2B)
SKYT = RGBColor(0x9F, 0xC0, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, size=18, bold=False, color=DARKTXT, align=PP_ALIGN.LEFT,
         first=False, bullet=False, space_after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = ('•  ' if bullet else '') + text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return p


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def band(s, t, h, color):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(t), SW, Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def header(s, title, kicker=None):
    band(s, 0, 1.15, NAVY)
    tf = box(s, 0.55, 0.14, 12.3, 0.98)
    if kicker:
        para(tf, kicker, size=12, bold=True, color=SKYT, first=True, space_after=2)
        para(tf, title, size=25, bold=True, color=WHITE)
    else:
        para(tf, title, size=27, bold=True, color=WHITE, first=True)


def add_table(s, rows, l, t, w, h, col_w=None, font=12, header_font=12,
              aligns=None):
    nrow = len(rows); ncol = len(rows[0])
    gf = s.shapes.add_table(nrow, ncol, Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = gf.table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(Inches(w).emu * cw / total))
    for r in range(nrow):
        for c in range(ncol):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.07); cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = (aligns[c] if aligns else PP_ALIGN.LEFT)
            run = p.add_run(); run.text = str(rows[r][c])
            run.font.name = 'Calibri'
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
                run.font.bold = True; run.font.size = Pt(header_font)
                run.font.color.rgb = WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = (LIGHT if r % 2 == 0 else WHITE)
                run.font.size = Pt(font); run.font.color.rgb = DARKTXT
    return tbl


def footer(s, n):
    tf = box(s, 11.9, 7.04, 1.3, 0.4)
    para(tf, str(n), size=11, color=GREY, align=PP_ALIGN.RIGHT, first=True)


# ---------------------------------------------------------- 1 TITLE
s = slide(); bg(s, NAVY)
band(s, 2.55, 0.035, BLUE)
tf = box(s, 0.9, 2.7, 11.5, 2.2)
para(tf, 'PETSE — TII 재투고 리뷰 대응 요약', size=40, bold=True,
     color=WHITE, first=True)
para(tf, 'From "just a wider margin" to continuous runtime re-verification',
     size=20, color=SKYT)
tf2 = box(s, 0.9, 5.0, 11.5, 1.7)
para(tf2, '핵심 기여 재정의: 마진 확장(X) → 승인된 작업의 실행 중 연속 재검증(O)',
     size=17, color=WHITE, first=True)
para(tf2, '신규 인과 실험 3종 + 리뷰어별 대응 매트릭스 + 논문 수치 전면 갱신',
     size=17, color=WHITE)
para(tf2, 'TII-26-3091 (Reject, 2026-07)  →  RA-L / RAS 재투고 준비',
     size=15, color=SKYT)

# ---------------------------------------------------------- 2 DIAGNOSIS
s = slide(); bg(s, WHITE)
header(s, '왜 부정적이었나 — 아이디어가 아니라 "안 보임"이 문제',
       '진단 (Root cause)')
rows = [
    ['리뷰어가 본 것', '왜 그렇게 읽혔나', '진짜 문제'],
    ['"그냥 마진을 넓힌 것"',
     '기여를 "불확실성 마진"으로 포장 → CBF/SSM과 차별화 안됨',
     '핵심(런타임 재검증)이 마진과 별개라는 증거가 논문에 없음'],
    ['LLM 버즈워드',
     '"LLM-enabled" 전면인데 정작 방법엔 LLM 없음 → 과대주장으로 읽힘',
     '신뢰도 전반 하락 → 다른 기여도 깎아서 봄'],
    ['완성도 부족',
     '통계 부재, 재현성(GitHub), 좋은 베이스라인, 하드웨어 누락',
     '완성도가 낮으면 참신성 주장에도 관대함을 줄임'],
]
add_table(s, rows, 0.55, 1.5, 12.25, 2.9, col_w=[3, 5, 5], font=15, header_font=15)
tf = box(s, 0.55, 4.75, 12.25, 2.1)
para(tf, '→ "아이디어가 틀렸다"가 아니라 "핵심 기여가 안 보인다"는 거절 — 받을 수 있는 거절 중 가장 좋은 종류.',
     size=17, bold=True, color=GREEN, first=True)
para(tf, '대응 전략: 핵심 기여를 리뷰어가 반박할 수 없는 형태(인과 실험)로 증명하고, 문서 전면에 재배치.',
     size=16, color=DARKTXT)
footer(s, 2)

# ---------------------------------------------------------- 3 REFRAME
s = slide(); bg(s, WHITE)
header(s, '재포지셔닝 — 마진이 아니라 런타임 재검증이 핵심', '전략')
tf = box(s, 0.55, 1.35, 12.25, 1.5)
para(tf, 'TOCTOU (Time-Of-Check ≠ Time-Of-Use): LLM/planner가 승인한 작업도 실행 중 문제가 생길 수 있다.',
     size=17, bold=True, color=NAVY, first=True)
para(tf, '승인은 취소 가능(revocable)이어야 한다 — PETSE는 매 사이클 승인된 작업을 재검증하고, 실제 상태가 제약을 위반하면 철회한다.',
     size=16, color=DARKTXT)
rows = [
    ['', '종전 (TII 투고본)', '재투고본 (재포지셔닝)'],
    ['제목', 'PETSE for LLM-enabled Mobile Robots',
     '... with Continuous Runtime Re-Verification (LLM 제거)'],
    ['기여 #1', 'Uncertainty-aware dynamic margin',
     'Revocable approval via continuous runtime re-verification'],
    ['마진의 위치', '핵심 주장', '재검증 임계값을 정하는 지원 역할(#2로 강등)'],
    ['핵심 증거', '없음', '마진 고정 2×2 — 연속 재검증이 공격면을 닫음 (100%→0%)'],
]
add_table(s, rows, 0.55, 3.1, 12.25, 2.9, col_w=[2.3, 4.5, 5.8], font=14, header_font=14)
footer(s, 3)


# ---------------------------------------------------------- 4-6 FIGURE SLIDES
def fig_slide(n, kicker, title, imgfile, bullets, caption):
    s = slide(); bg(s, WHITE)
    header(s, title, kicker)
    s.shapes.add_picture(os.path.join(FIG, imgfile), Inches(0.7), Inches(1.4),
                         width=Inches(8.4))
    tf = box(s, 9.35, 1.5, 3.65, 5.2)
    para(tf, '무엇을 보이나', size=15, bold=True, color=BLUE, first=True)
    for b in bullets:
        para(tf, b, size=13, bullet=True, space_after=9)
    tfc = box(s, 0.7, 6.4, 8.4, 0.85)
    para(tfc, caption, size=12, color=GREY, first=True)
    footer(s, n)


fig_slide(4, '신규 실험 ① — 인과 격리 (R1.1/R1.2/AE)',
          'TOCTOU 2×2 — 마진을 고정하고 모니터만 on/off',
          'toctou_ablation.png',
          ['마진 0.562m 고정, 모니터 OFF: 공격 10/10 진입',
           '모니터 ON: 0/10 — 100%→0% (McNemar p=1.95e-3)',
           '유일한 차이=런타임 모니터 → "그냥 마진" 반박',
           '경계 스윕이 해석 경계(1.389m)와 일치 → cherry-pick 아님'],
          '(a) 2×2 팩토리얼 · (b) 경계 스윕.   논문 fig:toctou / §Runtime Re-Verification')

fig_slide(5, '신규 실험 ② — 일반화 (R3.6)',
          '컨트롤러 일반화 — DWB vs RPP',
          'controller_generalization.png',
          ['런타임 모니터는 /cmd_vel에 붙어 컨트롤러-무관',
           '구조가 다른 두 Nav2 컨트롤러로 동일 100%→0%',
           'DWB(샘플러) · RPP(pure-pursuit), 속도 0.22m/s 매칭',
           '→ 특정 플래너 아티팩트가 아닌 실행계층 성질'],
          'DWB·RPP 모두 (OFF,above)=100%, 나머지 0%.   논문 fig:ctrlgen')

fig_slide(6, '신규 실험 ③ — 운영점 (연속 검사 비용)',
          '런타임 모니터 운영점 — 오발동 0, 적시 정지',
          'monitor_operating_point.png',
          ['양성 궤적 243개 → 오발동(nuisance) 0 (≤1.23%)',
           '55개 개입 전부 양의 여유로 정지, 위반 0/55',
           'S4(속도) median 3.1m · S5(TOCTOU) median 0.45m',
           '반응지연 median 50ms, per-cycle 147µs'],
          '(a) nuisance 0 · (b) 개입 시 여유거리 분포.   논문 fig:oppoint')

# ---------------------------------------------------------- 7 R1
s = slide(); bg(s, WHITE)
header(s, '리뷰어 1 대응 (독창성·엄밀성)', '리뷰어별 대응 매트릭스')
rows = [
    ['지적', '대응', '증거'],
    ['R1.1/1.2 "그냥 마진"', '마진 고정 2×2 → 모니터만으로 100%→0% + 재포지셔닝', 'fig:toctou'],
    ['R1.3/1.4 "보장 범위?"', 'COE(Certified Operating Envelope) 정의 + 3단계 주장', '§Method (예정)'],
    ['R1.5 "가정 붕괴"', 'TCB 신뢰표 + "둘 다 침해되면 보장X" 명시', '§Threat Model'],
    ['R1.6 "CBF 튜닝 불공정"', '(γ,δ) 30설정 스윕 → 최적도 33% 위반(구조적)', 'cbf_sensitivity.png'],
    ['R1 "additive 과보수"', 'additive≈p99.9; 적대적에서 RSS 100% 실패', 'margin_comparison.png'],
    ['R1 "fixed offset·covariance"', '지속 스푸핑 예산 Δ_spoof=0.462m 스윕', 'spoof_budget_sweep.png'],
    ['R1.8 "S1 하드웨어 누락"', '실기체 실험 예정 + 프롬프트 시도·모델 보고', '(진행 예정)'],
]
add_table(s, rows, 0.4, 1.4, 12.55, 5.4, col_w=[3.3, 6.4, 2.6], font=13, header_font=13)
footer(s, 7)

# ---------------------------------------------------------- 8 R2
s = slide(); bg(s, WHITE)
header(s, '리뷰어 2 대응 (재현성·표기)', '리뷰어별 대응 매트릭스')
rows = [
    ['지적', '대응', '증거'],
    ['R2 재현성 (코드 공개 없음)', 'GitHub 공개 저장소 + Setup에 링크(각주)', '§Setup 각주'],
    ['R2 표기 (ROS / AMCL)', 'ROS→ROS 2 전역 통일, 그림 내 철자 수정', '문서 전역'],
]
add_table(s, rows, 0.55, 1.5, 12.25, 1.8, col_w=[3.3, 6.1, 2.6], font=15, header_font=15)
tf = box(s, 0.55, 3.7, 12.25, 3.0)
para(tf, '성격: 논문 내용이 아니라 완성도·재현성 지적 — 방어보다 "정리"의 문제.',
     size=16, bold=True, color=NAVY, first=True)
para(tf, '코드·실험 하네스·결과 데이터를 저장소로 공개하고 Setup에 링크 → 재현성 확보.',
     size=15, bullet=True, color=DARKTXT)
para(tf, 'ROS 2 lifecycle node임을 본문 전역에서 명시(종전 "ROS"/"drop-in ROS node" 혼용 제거).',
     size=15, bullet=True, color=DARKTXT)
para(tf, '이런 완성도 지적을 남기면 참신성 주장에도 관대함이 줄어들므로 우선 정리.',
     size=15, bullet=True, color=DARKTXT)
footer(s, 8)

# ---------------------------------------------------------- 9 R3
s = slide(); bg(s, WHITE)
header(s, '리뷰어 3 대응 (일반화·실용성)', '리뷰어별 대응 매트릭스')
rows = [
    ['지적', '대응', '증거'],
    ['R3.5 좁은 복도 오거절', '거절구간도 footprint 여유 양수, ε로 조절', 'geometry_stress.png'],
    ['R3 인접구역 "마진 합산"', '합집합(union)이지 합산 아님 — 코드로 반박', 'geometry_stress.png'],
    ['R3 "ε 선택 지침 부재"', 'ε Pareto + 선택규칙(비용비→ε*)', 'epsilon_pareto.png'],
    ['R3.6 "플래너 교차검증"', 'DWB+RPP 2×2 동일 — 이번 세션 신규', 'controller_generalization.png'],
    ['R3.7 "fail-stop 복구불가"', 'degradation ladder, 위반 0 유지하며 가용성 복원', 'recovery_policies.png'],
    ['R3 실행계 문헌 누락', 'RTA/Simplex/RoboGuard 서베이 추가', '§Related Work (예정)'],
]
add_table(s, rows, 0.45, 1.5, 12.45, 4.8, col_w=[3.3, 6.3, 2.7], font=13.5, header_font=13.5)
tf = box(s, 0.55, 6.5, 12.25, 0.7)
para(tf, 'R3.6은 이번 세션에서 신규로 닫은 핵심 항목 — 실행계층 성질이 플래너와 무관함을 실증.',
     size=13, bold=True, color=GREEN, first=True)
footer(s, 9)

# ---------------------------------------------------------- 10 AE + stats
s = slide(); bg(s, WHITE)
header(s, 'AE · 통계 대응', '리뷰어별 대응 매트릭스')
rows = [
    ['지적', '대응', '증거'],
    ['AE 독창성 (R1.1 연장)', '런타임 재검증 재포지셔닝 + 인과 증명', 'fig:toctou, §reverify'],
    ['AE 마진 설계 정당화', 'additive vs RSS + ε Pareto', 'margin_comparison / epsilon_pareto'],
    ['AE CBF 튜닝 공정성', '(γ,δ) 30설정 스윕', 'cbf_sensitivity.png'],
    ['통계 부재 (공통)', 'Wilson CI, McNemar+Holm, Fisher, per-seed, rule-of-3',
     'Table detection/mcnemar'],
]
add_table(s, rows, 0.4, 1.4, 12.55, 2.9, col_w=[3.2, 6.0, 3.1], font=13, header_font=13)
tf = box(s, 0.55, 4.55, 12.25, 2.3)
para(tf, '핵심 통계 (S1+S3+S4+S5, seeds=20):', size=15, bold=True, color=NAVY, first=True)
para(tf, 'PETSE recall 100% [98.1,100], VR 0/300 (rule-of-3 ≤1.0%), FN=0', size=14, bullet=True)
para(tf, '모든 baseline 대비 McNemar p<0.001 (Holm), baseline 편향 discordant 쌍 0', size=14, bullet=True)
para(tf, 'margin-probe 재프레임: near/mid_boundary 거절=설계된 보수성 → FP=0', size=14, bullet=True)
footer(s, 10)

# ---------------------------------------------------------- 10 numbers
s = slide(); bg(s, WHITE)
header(s, '논문 수치 전면 갱신 (구 → 신)', '수치 일관성')
rows = [
    ['항목', '종전 (구)', '재투고본 (신)'],
    ['시험 규모', '1,920 trials', 'seeds=20, baseline 300/method, 총 2,584'],
    ['PETSE F1', '0.917', '1.000 [1.000,1.000] (margin-probe 재프레임)'],
    ['PETSE FP', '40', '0 (near/mid_boundary 별도 보고)'],
    ['최강 baseline', 'CBF-Adaptive F1 0.548', 'CBF-Adaptive F1 0.592, recall 0.420'],
    ['baseline 수', '5', '6 (RoboGuard 추가)'],
    ['유의성', '없음', 'McNemar/Holm p<0.001, CI, per-seed'],
    ['새 그림', '—', '2×2·컨트롤러·운영점 + RSS/CBF/ε/기하/스푸핑'],
]
add_table(s, rows, 0.55, 1.4, 12.25, 4.6, col_w=[2.7, 4.2, 6.0], font=13.5, header_font=14)
tf = box(s, 0.55, 6.2, 12.25, 0.9)
para(tf, '→ 초록·기여문·본문·표가 모두 새 수치로 일치 (구 수치 잔존 0 확인).',
     size=15, bold=True, color=GREEN, first=True)
footer(s, 11)

# ---------------------------------------------------------- 11 roadmap
s = slide(); bg(s, WHITE)
header(s, '진행 현황 · 남은 작업', '로드맵')
tf = box(s, 0.55, 1.4, 6.0, 5.4)
para(tf, '✅ 완료', size=17, bold=True, color=GREEN, first=True)
for t in ['핵심 인과 실험 3종 (2×2·컨트롤러·운영점)',
          '분석·시뮬 실험 (CBF·RSS·ε·기하·스푸핑·복구)',
          '통계 (CI·McNemar·Holm·per-seed)',
          '논문: 제목·초록·기여문·수치 표 전면 갱신',
          '신규 결과 소절 + 그림 3장 삽입']:
    para(tf, t, size=13.5, bullet=True, color=DARKTXT, space_after=8)
tf2 = box(s, 6.85, 1.4, 6.0, 5.4)
para(tf2, '⏳ 남음', size=17, bold=True, color=RED, first=True)
for t in ['본문: COE 정의 + TCB 표 (Method/Threat)',
          'Limitations 재작성 (자기모순 제거, fail-stop→recovery)',
          'Related Work: 실행계 문헌 서베이+bib',
          'S1 하드웨어 실기체 실험 (선택)',
          'GitHub 저장소 공개·README',
          'Overleaf 컴파일 확인 (로컬 LaTeX 없음)']:
    para(tf2, t, size=13.5, bullet=True, color=DARKTXT, space_after=8)
footer(s, 12)

prs.save(OUT)
print('saved', OUT, '| slides:', len(prs.slides._sldIdLst))
