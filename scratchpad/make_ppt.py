#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""PETSE revision-response deck for the advisor: reviewer critique → experiment → result."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FIG = "/home/jim/ros2_motion_planning_tutorials/experiment_results/gazebo_s1_s6"
SCR = "/home/jim/ros2_motion_planning_tutorials/scratchpad"
KFONT = "Malgun Gothic"   # Korean-capable; viewers substitute if absent

INK   = RGBColor(0x1a, 0x2a, 0x33)
TEAL  = RGBColor(0x2a, 0x9d, 0x8f)
RED   = RGBColor(0xc0, 0x39, 0x2b)
AMBER = RGBColor(0xd6, 0x8a, 0x10)
GREY  = RGBColor(0x55, 0x5f, 0x66)
LIGHT = RGBColor(0xf3, 0xf6, 0xf7)
WHITE = RGBColor(0xff, 0xff, 0xff)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set(tf_or_para, text, size, color=INK, bold=False, font=KFONT, align=None):
    p = tf_or_para
    p.text = text
    for r in p.runs:
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
        r.font.color.rgb = color
    if align is not None:
        p.alignment = align


def band(slide, color, y, h):
    box = slide.shapes.add_shape(1, 0, y, SW, h)
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.fill.background()
    box.shadow.inherit = False
    return box


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def title_bar(slide, tag, title, tagcolor=TEAL):
    band(slide, LIGHT, 0, Inches(1.15))
    strip = slide.shapes.add_shape(1, 0, 0, Inches(0.18), Inches(1.15))
    strip.fill.solid(); strip.fill.fore_color.rgb = tagcolor; strip.line.fill.background()
    strip.shadow.inherit = False
    tf = add_textbox(slide, Inches(0.5), Inches(0.12), Inches(12.4), Inches(0.44))
    _set(tf.paragraphs[0], tag, 13, tagcolor, bold=True)
    tf2 = add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12.4), Inches(0.6))
    _set(tf2.paragraphs[0], title, 25, INK, bold=True)


def bullets(slide, x, y, w, h, items, size=15, gap=6):
    tf = add_textbox(slide, x, y, w, h)
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        # two accepted forms:
        #   (level:int, text:str[, color[, bold]])         -> indented sub-bullet
        #   (text:str[, color[, bold]])  or  plain str      -> level-0 header line
        if isinstance(it, tuple) and isinstance(it[0], int):
            lvl, txt = it[0], it[1]
            color = it[2] if len(it) > 2 else INK
            bold = it[3] if len(it) > 3 else False
        elif isinstance(it, tuple):
            lvl, txt = 0, it[0]
            color = it[1] if len(it) > 1 else INK
            bold = it[2] if len(it) > 2 else False
        else:
            lvl, txt, color, bold = 0, it, INK, False
        p.level = lvl
        mark = "•  " if lvl == 0 else "–  "
        _set(p, mark + txt, size - lvl, color, bold=bold)
        p.space_after = Pt(gap)
    return tf


def add_image_fit(slide, path, x, y, maxw, maxh):
    if not os.path.exists(path):
        tf = add_textbox(slide, x, y, maxw, Inches(0.4))
        _set(tf.paragraphs[0], f"[missing: {os.path.basename(path)}]", 11, RED)
        return
    from PIL import Image
    iw, ih = Image.open(path).size
    r = min(maxw / iw, maxh / ih)
    w, h = int(iw * r), int(ih * r)
    slide.shapes.add_picture(path, x + (maxw - w) // 2, y + (maxh - h) // 2, width=w, height=h)


def takeaway(slide, text, color=TEAL):
    y = SH - Inches(0.95)
    bar = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.33), Inches(0.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    _set(tf.paragraphs[0], "결론:  " + text, 14, WHITE, bold=True)


# ───────────────────────────────────────────── 1. TITLE
s = prs.slides.add_slide(BLANK)
band(s, INK, 0, SH)
band(s, TEAL, Inches(3.05), Inches(0.06))
tf = add_textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.6))
_set(tf.paragraphs[0], "PETSE 논문 개정 대응", 40, WHITE, bold=True)
p = tf.add_paragraph(); _set(p, "이동로봇 실행 계층 방어: 센서 스푸핑·TOCTOU 공격에 대한 지속 런타임 재검증", 18, TEAL, bold=True)
tf2 = add_textbox(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(2.6))
for i, t in enumerate([
    "IEEE TII (TII-26-3091) 리젝트 (2026-07) → NDSS(보안 학회) 제출 준비",
    "핵심 질문: 리뷰어 지적을 실험으로 어떻게 대응했는가",
    "구성:  ① 리뷰어 6대 지적  →  ② 대응 실험 설계  →  ③ 결과",
]):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    _set(p, t, 17, WHITE); p.space_after = Pt(10)
tf3 = add_textbox(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6))
_set(tf3.paragraphs[0], "조준영 · 지도교수 보고용", 14, GREY)

# ───────────────────────────────────────────── R. REVIEWER COMMENTS (verbatim)
s = prs.slides.add_slide(BLANK)
title_bar(s, "받은 리뷰 (TII)", "리뷰어 1 — 신규성·엄밀성·위협모델 전반 지적 (리젝트 주도)", tagcolor=RED)
bullets(s, Inches(0.55), Inches(1.35), Inches(12.3), Inches(5.5), [
    (1, "① 신규성 제한 — geofencing·불확실성 팽창·braking·런타임 모니터링의 조합, 근본적 새 프레임워크 아님"),
    (1, "② additive safety margin이 보수적 — tracking error·localization·latency 결합효과를 이중계산할 수 있음"),
    (1, "③ 결정론적 보장이 bounded-disturbance 가정(latency·velocity·deceleration·covariance)에 의존 — 실제 공격 시 미충족 가능"),
    (1, "④ fail-stop은 fallback일 뿐 — 적대적 가정위반 하에서 견고한 안전 미제공"),
    (1, "⑤ 사이버-물리 보안 약함 — 저수준 액추에이터·펌웨어 공격·PETSE 노드 자체 장악을 위협모델에 미포함"),
    (1, "⑥ baseline 불공정 — CBF는 튜닝에 매우 민감, 파라미터 민감도·최적 보정 미제시"),
    (1, "⑦ 실기 검증 제한 — 소수 trial·저속 단순 플랫폼·단순 금지구역 기하"),
    (1, "⑧ 가장 중요한 위험목표 시나리오(S1)가 실기 하드웨어에서 미시연"),
    (1, "⑨ 통계 엄밀성 부족 — 신뢰구간·유의성 검정·시드 간 분산 분석 없음"),
    (1, "⑩ LLM은 대체로 동기부여용 — 일반 실행계층 공간 안전 모니터일 뿐 LLM 특화 기법 아님"),
], size=13, gap=8)

s = prs.slides.add_slide(BLANK)
title_bar(s, "받은 리뷰 (TII)", "리뷰어 2 (긍정) · 리뷰어 3 (건설적)", tagcolor=AMBER)
tf = add_textbox(s, Inches(0.55), Inches(1.3), Inches(12.3), Inches(0.4))
_set(tf.paragraphs[0], "리뷰어 2 — 전반적으로 호의적", 15, TEAL, bold=True)
bullets(s, Inches(0.7), Inches(1.75), Inches(12.1), Inches(1.6), [
    (1, "시뮬레이션+실기 검증, SoTA(SELP·CBF·SSM) 비교와 상세 ablation은 고무적, 잘 쓰인 논문"),
    (1, "일부 문장 더 자연스럽게 재작성 · 표기/약어 일관성(예: 'ACML'→AMCL 오타) · ROS vs ROS2 혼용 정리"),
    (1, "셋업·실험을 담은 GitHub 페이지 기대"),
], size=13, gap=5)
tf = add_textbox(s, Inches(0.55), Inches(3.35), Inches(12.3), Inches(0.4))
_set(tf.paragraphs[0], "리뷰어 3 — 실행계층 baseline 부재 지적 + 6개 구체 질문", 15, TEAL, bold=True)
bullets(s, Inches(0.7), Inches(3.8), Inches(12.1), Inches(3.0), [
    (1, "(총평) 실행계층 방어 baseline 미인용 — 적응형 기하 모니터를 baseline으로 구현해 additive 불확실성 마진의 이점을 분리 권장"),
    (1, "1. 운영자가 위험 파라미터 ε를 보수성 vs 가용성 사이에서 어떻게 선택?"),
    (1, "2. S1: LLM이 위험목표 생성 거부 — 시도한 프롬프트 수·수동 대체목표 사용 여부 보고"),
    (1, "3. Algorithm 1이 팽창 zone에 걸치는 경로도 거부 → 좁은 통로에서 잦은 오탐 우려, 명확화"),
    (1, "4. 정적·고립 폴리곤 zone만; 다중 인접/시변 zone; 두 팽창 zone 중첩 시 마진합산이 차단 유발?"),
    (1, "5. eq(3) envelope의 fitting 절차·다중 컨트롤러(DWB vs pure pursuit)·지형 검증 없음"),
    (1, "6. fail-stop이 로봇을 통로에 방치·복구 불가 상태로 만들 수 있음 — 해결책은?"),
], size=12.5, gap=5)

# ───────────────────────────────────────────── R2. REVIEW → RESPONSE MAP
s = prs.slides.add_slide(BLANK)
title_bar(s, "리뷰 → 대응 매핑", "실제 지적별 대응 (대부분 이번 개정에서 실험 완료)", tagcolor=INK)
map_rows = [
    ("R1-⑨ 통계(CI·유의성·시드분산)", "20시드·Wilson CI·McNemar+Holm·시드별 분산", "완료"),
    ("R1-⑥/R3 CBF 공정성·튜닝 민감도", "CBF-Adaptive baseline + 파라미터 민감도 sweep", "완료"),
    ("R1-⑤ 위협모델(액추에이터·펌웨어·PETSE노드)", "TCB·fail-closed 논거 + A1 명시 (금일 보강)", "완료"),
    ("R1-④/R3-6 fail-stop 견고성·복구", "fail-closed + safe-hold 자동복구 실험", "완료"),
    ("R1-② additive 마진 이중계산", "RSS vs additive Monte-Carlo (보수성 33% 정량화)", "완료"),
    ("R1-⑩/R2 LLM 위치", "LLM-enabled 제거 → 일반 실행계층으로 재포지셔닝", "완료"),
    ("R3-1 ε 선택 가이드", "교차채널 탐지기 ROC-AUC + 동작점 제시", "완료"),
    ("R3-3 좁은 통로 오탐", "좁은 통로 멀티시드·멀티방법 실험 (금일)", "완료"),
    ("R3-4 인접·중첩 zone / 마진 중첩", "기하 sweep(multi-zone) + 중첩 마진 (금일)", "완료"),
    ("R3-5 컨트롤러 일반화(DWB vs RPP)", "envelope 컨트롤러-불변성 검증", "완료"),
    ("R3-총평 실행계층 baseline", "RoboGuard(승인계층) 추가", "완료"),
    ("R1-⑧·⑦ / R3-2 S1 실기·프롬프트 수", "실기 하드웨어 검증 + 프롬프트 시도수 보고", "예정"),
    ("R2 GitHub·표기·ROS2 일관성", "오픈소스 공개 + 표기/약어 정리", "예정"),
]
top = Inches(1.35); rowh = Inches(0.415)
xs = [Inches(0.5), Inches(6.0), Inches(11.9)]; ws = [Inches(5.4), Inches(5.8), Inches(1.0)]
hb = band(s, INK, top, Inches(0.38))
for j, htxt in enumerate(["리뷰어 지적", "대응 실험 / 조치", "상태"]):
    tf = add_textbox(s, xs[j], top + Inches(0.02), ws[j], Inches(0.34), MSO_ANCHOR.MIDDLE)
    _set(tf.paragraphs[0], htxt, 12, WHITE, bold=True)
for i, (crit, exp, st) in enumerate(map_rows):
    yy = top + Inches(0.38) + rowh * i
    if i % 2 == 0:
        band(s, LIGHT, yy, rowh)
    for j, txt in enumerate([crit, exp, st]):
        tf = add_textbox(s, xs[j], yy + Inches(0.02), ws[j], rowh - Inches(0.04), MSO_ANCHOR.MIDDLE)
        col = TEAL if (j == 2 and st == "완료") else (AMBER if j == 2 else INK)
        _set(tf.paragraphs[0], txt, 10.5, col, bold=(j == 2))

# ───────────────────────────────────────────── 2. BACKGROUND
s = prs.slides.add_slide(BLANK)
title_bar(s, "배경", "PETSE란 무엇이고, 왜 개정이 필요한가")
bullets(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(3.0), [
    ("PETSE = 이동로봇용 '실행 계층(execution-layer) geofence' — 목표 승인 후에도 실행 중 지속적으로 재검증하여 금지구역 진입을 원천 차단", INK, True),
    (1, "기존 방법: 계획/승인 시점에만 검사 → 승인 후(TOCTOU)·스푸핑 공격에 무방비"),
    (1, "PETSE: 실행 시점에 위치·경로·교차채널(AMCL vs odom)을 재검증 → 승인을 실시간 취소(fail-stop)"),
    ("TII 리젝트 (2026-07): 아이디어는 인정하나 '실험적 엄밀성·현실성' 지적 → NDSS(보안) 목표로 위협·방어 중심 재프레이밍 + 실험 보강", RED, True),
])
band(s, LIGHT, Inches(4.8), Inches(1.7))
tf = add_textbox(s, Inches(0.6), Inches(4.95), Inches(12.1), Inches(1.4))
_set(tf.paragraphs[0], "위협 모델 (4개 시나리오)", 16, TEAL, bold=True)
for t in ["S1 금지구역 목표  ·  S2 경로 통과  ·  S3 살라미(점진)  ·  S4 스푸핑/TOCTOU 공격",
          "가정: odometry 채널 + 모니터 연산은 비탈취(A1) — 그 이상은 물리 HW 안전채널(범위 밖)"]:
    p = tf.add_paragraph(); _set(p, t, 14, INK); p.space_after = Pt(4)
takeaway(s, "실행 계층 재검증이 핵심 기여 — 리뷰어 지적은 이를 '얼마나 엄밀·현실적으로 입증했나'에 집중")

# helper for figure slides
def figure_slide(tag, title, fig_path, bullet_items, take, tagcolor=TEAL, imgw=7.7):
    s = prs.slides.add_slide(BLANK)
    title_bar(s, tag, title, tagcolor=tagcolor)
    bullets(s, Inches(0.55), Inches(1.4), Inches(4.9), Inches(4.6), bullet_items, size=13.5, gap=7)
    add_image_fit(s, fig_path, Inches(5.5), Inches(1.35), Inches(imgw), Inches(4.9))
    takeaway(s, take, tagcolor)
    return s

# ───────────────────────────────────────────── 4. ① STATISTICS
figure_slide("① 통계적 엄밀성", "신뢰구간·유의성 검정 추가",
    os.path.join(FIG, "money_2x2/final_unified_2x2.png"),
    [("실험 설계", TEAL, True),
     (1, "seeds 0–19 (20 반복), 총 2,794 trial"),
     (1, "Wilson 95% CI, rule-of-three, McNemar 쌍대검정 + Holm-Bonferroni 보정"),
     ("결과", TEAL, True),
     (1, "PETSE recall 100% [98.1, 100], FN=0"),
     (1, "위반율(VR) 0/301 (rule-of-3 ≤ 1.0%)"),
     (1, "F1 0.910 [0.880, 0.935]"),
     (1, "최고 baseline (CBF-Adaptive) F1 0.519, recall 42%"),
     (1, "모든 McNemar p < 0.001 (Holm 후), 역전쌍 0"),
    ],
    "통계적으로 PETSE 우위가 유의 — CI·유의성 지적 해소")

# ───────────────────────────────────────────── 5. ADAPTIVE ATTACK (NDSS headline)
figure_slide("적응형 공격 (NDSS 핵심)", "방어를 아는 공격자: 협조 이중채널 스푸핑",
    os.path.join(FIG, "money_2x2/coord_epsilon_curve.png"),
    [("왜 핵심인가", RED, True),
     (1, "NDSS는 '방어 메커니즘을 아는 적응형 공격자'를 반드시 요구"),
     (1, "PETSE의 방어 = 교차채널 감지(AMCL vs odom 불일치)"),
     ("적응형 공격", AMBER, True),
     (1, "감지를 무력화하려 LiDAR+odom을 동시 스푸핑"),
     (1, "두 채널을 함께 위조해 교차채널 잔차를 목표 ε로 유지"),
     ("결과 (ε 스윕)", TEAL, True),
     (1, "잔차가 τc=0.95m 넘으면 침해 0 → 우아한 성능 저하"),
     (1, "실시간 협조 불완전(잔차 +0.2~0.3m) → 회피엔 ε≲0.65 필요"),
     (1, "= 로봇 온보드+odom까지 완전 장악(A1 범위 밖)이 전제"),
    ],
    "적응형 공격도 온보드 완전 장악 없이는 회피 불가 — '이미 이긴 공격자'만 성공", tagcolor=RED)

# ───────────────────────────────────────────── 5b. ATTACK CAPABILITY LADDER (NDSS threat model)
figure_slide("공격자 능력 (NDSS 위협모델)", "각 공격에 필요한 장악 수준 — PETSE 방어 범위",
    os.path.join(SCR, "attack_ladder.png"),
    [("공격 사다리 (하 → 상)", RED, True),
     (1, "S1/S2 목표 주입: 명령 입력만 (LLM/오퍼레이터 채널)"),
     (1, "S3 파라미터 변조: 네트워크 발판(설정 접근)"),
     (1, "S4 cmd_vel 주입: ROS/DDS 퍼블리시 (동일 도메인)"),
     (1, "LiDAR 스푸핑: 물리적 가시선(SW 침투 불필요)"),
     ("실증 (§III)", TEAL, True),
     (1, "ROS2 기본 Fast-DDS = 무인증 → 임의 프로세스가"),
     (1, "  /cmd_vel 주입 가능 (도메인 격리 실험서 확인)"),
     ("범위 밖 (A1)", AMBER, True),
     (1, "협조 LiDAR+odom = 온보드 완전 장악 (기능안전 영역)"),
    ],
    "PETSE는 1~4단계(명령·네트워크·DDS·물리 스푸핑)를 실행계층에서 차단 — 온보드 완전장악만 범위 밖", tagcolor=RED)

# ───────────────────────────────────────────── 6a. ④ NARROW — SETUP
s = prs.slides.add_slide(BLANK)
title_bar(s, "④ 좁은 통로", "실험 설정: 양측 대칭 금지구역으로 안전폭 스윕", tagcolor=AMBER)
add_image_fit(s, os.path.join(SCR, "narrow_setup.png"), Inches(0.4), Inches(1.35), Inches(12.5), Inches(4.0))
bullets(s, Inches(0.8), Inches(5.35), Inches(11.7), Inches(1.1), [
    (1, "y=±h 대칭 금지구역이 y=0 통로를 위·아래로 압박 → 안전폭 = 2(h − M), 마진 M≈0.55 m"),
    (1, "안전폭을 1.2 m(안전) → 0.6 m(경계) → 0.0 m(통과불가)로 조절, 각 3방법 × 5시드"),
], size=13, gap=5)
takeaway(s, "안전~통과불가를 연속으로 스윕해 '과잉차단 vs 올바른 차단'을 동시에 검증", AMBER)

# ───────────────────────────────────────────── 6b. ④ NARROW — RESULT
figure_slide("④ 좁은 통로", "좁은 통로 결과 (멀티시드·멀티방법)",
    os.path.join(FIG, "narrow/narrow_corridor.png"),
    [("리뷰어 우려", AMBER, True),
     (1, "정지시켜 길을 막아 정상 통행을 방해하지 않나?"),
     ("결과", TEAL, True),
     (1, "안전폭 1.2m: PETSE 5/5 통과 (과잉차단 0)"),
     (1, "통과불가(마진 겹침): PETSE 5/5 정지 (유일)"),
     (1, "동일 0.55m 마진의 CBF는 point-goal이라 5/5 밀고 통과"),
    ],
    "판별자는 마진 크기가 아니라 경로·실행 인지 — 안전은 통과, 불가는 차단", tagcolor=AMBER)

# ───────────────────────────────────────────── 7. ④ GENERALIZATION
figure_slide("④ 일반화", "금지구역 기하 일반화 sweep",
    os.path.join(FIG, "geom/geom_generalization.png"),
    [("실험 설계", TEAL, True),
     (1, "동일 맵에서 금지구역 기하 변형"),
     (1, "compact / shifted / wide-thin / multi-zone × 3 seeds"),
     ("결과", TEAL, True),
     (1, "no_guard: 모든 기하에서 12/12 침해"),
     (1, "PETSE: 모든 기하에서 0/12 침해"),
     (1, "PETSE의 0% 위반율이 기하에 무관하게 일반화"),
    ],
    "특정 형상에 과적합 아님 — 다양한 금지구역 기하에서 일관되게 차단")

# ───────────────────────────────────────────── 8. ④ REALISM FoV
figure_slide("④ 현실성", "물리적으로 현실적인 제한 시야(FoV) 스푸핑",
    os.path.join(FIG, "money_2x2/realism_fov_curve.png"),
    [("리뷰어 우려", AMBER, True),
     (1, "이상적 360° 전면 스캔 교체는 비현실적"),
     ("실험 설계", TEAL, True),
     (1, "스푸퍼 시야를 360→180→90→45°로 제한 (×3 seeds)"),
     ("결과", TEAL, True),
     (1, "FoV↓ → AMCL 변위 단조 붕괴 (d_abs 1.00→0.30)"),
     (1, "≤180°: 실제 빔이 AMCL을 고정 → 공격 무력화"),
     (1, "모든 FoV에서 침해 0/3 (강할 땐 PETSE가 탐지)"),
    ],
    "공격 강도는 이상화 산물 — 현실적 스푸핑은 AMCL 다중빔 융합이 방어", tagcolor=AMBER)

# ───────────────────────────────────────────── 9. ⑤ ROC
figure_slide("⑤ 탐지기 동작점", "교차채널 탐지기 ROC · ε 가이드",
    os.path.join(FIG, "money_2x2/detector_roc.png"),
    [("리뷰어 우려", AMBER, True),
     (1, "임계값(ε/τc) 선택 근거·동작점 제시 필요"),
     ("실험 설계", TEAL, True),
     (1, "AMCL-odom 교차채널 잔차 기반 탐지기 ROC 산출"),
     (1, "CUSUM(누적) vs memoryless(순간) 비교"),
     ("결과", TEAL, True),
     (1, "ROC-AUC로 동작점 정량화 → τc=0.95m 근거 제시"),
     (1, "청정 오탐 없이 스푸핑 분리 가능한 마진 확인"),
    ],
    "임계값 선택을 ROC로 정당화 — 동작점 가이드 제공")

# ───────────────────────────────────────────── 10. SPOOFING HIJACK (warehouse) — headline security result
figure_slide("핵심 보안 결과", "스푸핑 카메라-하이잭 — 계획 계층은 전부 뚫림, PETSE만 차단",
    os.path.join(SCR, "fig_wh_hijack.png"),
    [("위협 (NDSS 핵심)", RED, True),
     (1, "탈취된 로봇이 카메라로 기밀 구역 촬영 시도"),
     (1, "map-consistent 스푸핑이 AMCL을 +Y로 위조 →"),
     (1, "진짜 로봇을 −Y 금지구역으로 유도 (승인 후 TOCTOU)"),
     ("baseline = 계획 계층", TEAL, True),
     (1, "CBF/SSM: 실행 시점 재검증 없음 (논문 논지)"),
     ("결과 (5시드 × 4방법, warehouse)", TEAL, True),
     (1, "No Guard·SELP·CBF: 5/5 구역 침입(FN)", RED, True),
     (1, "PETSE: 0/5 — 교차채널 런타임 재검증 fail-stop", TEAL, True),
    ],
    "실행 계층 재검증이 승인 후 스푸핑 공격을 유일하게 차단 — NDSS 핵심 보안 기여", tagcolor=RED, imgw=7.4)

# ───────────────────────────────────────────── 10b. ASSUMPTION VIOLATION (R1-③, velocity bound)
figure_slide("① 가정 위반 (R1-③)", "선언된 속도 상한을 넘기는 공격 — 실행시점 재검증의 필요성",
    os.path.join(SCR, "fig_velviol.png"),
    [("리뷰어 지적 (R1-③)", AMBER, True),
     (1, "마진이 v_max·τ 가정에 의존 → 가정 위반 시 무력화?"),
     ("공격", RED, True),
     (1, "cmd_vel 주입으로 과속 (선언 v_max=0.5 → 실주행 1.5 m/s, 3×)"),
     ("베이스라인: 반응형 고정마진", RED, True),
     (1, "선언 v_max 신뢰 → 고정 0.55m 버퍼 (전방예측 없음)"),
     (1, "제동거리(½v²/a) > 0.55m → 오버슈트하여 마진 침해"),
     ("PETSE 대응", TEAL, True),
     (1, "실행시점 재검증 + 실측 속도(95th pctile)로 마진 재계산"),
     ("결과 (각 5시드, 이격거리)", TEAL, True),
     (1, "반응형 고정마진: 안전마진(0.55m) 5/5 침해 (이격 ≤0.29m)", RED, True),
     (1, "  → 로봇 몸체가 이미 금지구역 침범, 중심진입 2/5", RED, False),
     (1, "PETSE: 0/5 침해, 2.2m 이격 유지", TEAL, True),
    ],
    "선언 v_max 신뢰 고정마진은 과속에 자기 안전마진을 5/5 전부 침해 — PETSE만 이격 유지", tagcolor=AMBER, imgw=7.4)

# ───────────────────────────────────────────── 10c. TIME-VARYING ZONE (R3-4)
figure_slide("③ 시변 금지구역 (R3-4)", "주행 중 활성화되는 금지구역 — 승인 시점 방식의 한계",
    os.path.join(SCR, "fig_tvzone.png"),
    [("리뷰어 지적 (R3-4)", AMBER, True),
     (1, "금지구역이 시간에 따라 변하면? (동적 안전영역)"),
     ("시나리오", RED, True),
     (1, "목표 (6.5,0)를 빈 경로에서 승인 → 주행 시작"),
     (1, "로봇이 x∈[0.5,2.8] 통과 시 외부에서 zone 활성화"),
     (1, "  (/petse/inject_zone, 승인 이후 = TOCTOU)"),
     ("결과 (각 5시드)", TEAL, True),
     (1, "No Guard·SELP: 승인 시점만 검사 → 5/5 진입", RED, True),
     (1, "PETSE: 연속 재검증으로 활성화 즉시 fail-stop 0/5", TEAL, True),
     (1, "(스푸핑 없는 시나리오 → 가드형 CBF도 0/5 방어;", INK, False),
     (1, " 하이잭과 달리 교차채널 감지 불필요)", INK, False),
    ],
    "승인 시점 검사만으론 시변 금지구역 방어 불가 — 연속 재검증(PETSE)이 핵심", tagcolor=AMBER, imgw=7.4)

# ───────────────────────────────────────────── 11b. WHY NOT JUST TURN IT OFF (TCB / fail-closed)
s = prs.slides.add_slide(BLANK)
title_bar(s, "심화 질문 (예상 리뷰)", "\"이미 침투했으면 PETSE를 끄면 되지 않나?\"", tagcolor=RED)
add_image_fit(s, os.path.join(SCR, "tcb_failclosed.png"), Inches(0.35), Inches(1.25), Inches(7.4), Inches(4.7))
bullets(s, Inches(7.95), Inches(1.4), Inches(5.1), Inches(4.6), [
    ("판별선 = TCB 장악 여부", RED, True),
    (1, "'로봇 침투' = 보통 내비 스택 장악 ≠ PETSE 코어"),
    (1, "PETSE는 분리된 작은 신뢰 코어로 감시"),
    ("끄기는 우회가 아님 (fail-closed)", TEAL, True),
    (1, "가드가 모터 토픽의 유일 게이트웨이"),
    (1, "끄면 → cmd_vel_safe 끊김 → 로봇 정지"),
    (1, "워치독: PETSE 무효 → 구동 차단 (dead-man)"),
    (1, "진입하려면 켜둔 채 속여야 → 교차채널 감지"),
    ("정직한 범위 밖 (A1)", AMBER, True),
    (1, "안전 컨트롤러 물리 장악 / 모터 재배선"),
    (1, "기능안전 표준(PLC·light curtain) 영역"),
], size=12.5, gap=6)
takeaway(s, "끄면 정지 · 진입하려면 켜둔 채 속여야 함 → '그냥 끄기'는 자멸적 (논문 §III A1에 fail-closed 논거 보강)", RED)

# ───────────────────────────────────────────── 12. SUMMARY
s = prs.slides.add_slide(BLANK)
title_bar(s, "종합", "결과 요약 — 모든 지적에 실험적 대응", tagcolor=INK)
bullets(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(4.4), [
    ("정량 헤드라인 (20 seeds)", TEAL, True),
    (1, "PETSE recall 100% · FN 0 · VR 0/301 · F1 0.910  vs  최고 baseline F1 0.519"),
    (1, "모든 McNemar p<0.001 (Holm), 역전쌍 0"),
    ("강건성·현실성", TEAL, True),
    (1, "좁은 통로: 안전 통과 + 통과불가 차단 (판별자=경로·실행 인지)"),
    (1, "기하 일반화 0/12 · 제한 FoV 스푸핑 침해 0/모든 FoV · 협조공격 ε≲0.65만 회피(범위 밖)"),
    (1, "스푸핑 하이잭(warehouse, 5시드): 계획 계층(No Guard·SELP·CBF) 5/5 침입, PETSE 0/5 (유일 차단)"),
    ("NDSS 공격 확장", RED, True),
    (1, "공격자 능력 사다리 + DDS 무인증 /cmd_vel 주입 실증(§III)"),
    (1, "가정 위반(R1-③, 속도상한 초과): 반응형 고정마진 안전마진 5/5 침해, PETSE 0/5(2.2m 이격) · 시변 금지구역(R3-4): 승인시점 방식 5/5 진입, PETSE 0/5"),
    ("남은 과제", AMBER, True),
    (1, "S1 실기 하드웨어 검증 (현재 Gazebo 고정밀 재현으로 대체)"),
], size=14, gap=8)
takeaway(s, "TII 지적 대부분 실험 완료 · S1 실기 예정 → NDSS 제출 준비", TEAL)

# ───────────────────────────────────────────── 13. CLOSING
s = prs.slides.add_slide(BLANK)
band(s, INK, 0, SH)
band(s, TEAL, Inches(3.6), Inches(0.06))
tf = add_textbox(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.4))
_set(tf.paragraphs[0], "핵심 메시지", 30, TEAL, bold=True)
tf2 = add_textbox(s, Inches(0.9), Inches(3.4), Inches(11.5), Inches(2.6))
for i, t in enumerate([
    "PETSE의 기여 = 실행 계층의 지속적 재검증 (승인 후 공격을 실시간 취소)",
    "리뷰어 지적을 통계·baseline·위협모델·현실성·동작점 실험으로 정면 대응",
    "원래 warehouse 맵에서 '정상 통행 보존 + 스푸핑/TOCTOU 공격 차단'을 입증",
]):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    _set(p, "•  " + t, 17, WHITE); p.space_after = Pt(12)

out = "/home/jim/ros2_motion_planning_tutorials/PETSE_revision_response.pptx"
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
